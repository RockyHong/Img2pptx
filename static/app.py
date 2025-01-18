from flask import Flask, render_template, request, send_from_directory
import os
from pptx import Presentation
from PIL import Image

# Initialize the Flask app
app = Flask(__name__)

# Define the route for the home page
@app.route("/", methods=["GET", "POST"])
def home():
    if request.method == "POST":
        # Process the form data
        print("Form submitted!")
        return "Form processed successfully!"
    # Serve the HTML page for GET requests
    
    """
    Handle the uploaded ZIP file, extract images, and create a PowerPoint
    presentation based on user input for dimensions and fit mode.
    """
    # Extract form data from the HTML form
    slide_width_px = int(request.form.get("slide_width_px", 1920))
    slide_height_px = int(request.form.get("slide_height_px", 1080))
    fit_mode = request.form.get("fit_mode", "fit_width")

    # Access the uploaded ZIP file
    images_zip = request.files.get("images_zip")

    if not images_zip or not images_zip.filename.endswith(".zip"):
        return "Error: Please upload a valid .zip file containing images."

    # Create a temporary directory to store extracted images
    temp_dir = "temp_upload"
    os.makedirs(temp_dir, exist_ok=True)

    # Extract the uploaded ZIP file to the temporary directory
    zip_path = os.path.join(temp_dir, images_zip.filename)
    images_zip.save(zip_path)
    os.system(f"unzip -o {zip_path} -d {temp_dir}")  # Extract the ZIP contents

    # Create a new PowerPoint presentation
    presentation = Presentation()

    for root, _, files in os.walk(temp_dir):
        for file in files:
            if file.lower().endswith((".png", ".jpg", ".jpeg")):
                # Add a new slide
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])

                # Get the image path
                img_path = os.path.join(root, file)

                # Open the image to get its dimensions
                img = Image.open(img_path)
                img_width, img_height = img.size

                # Adjust image size and position based on the fit mode
                if fit_mode == "fit_width":
                    scaled_height = int(slide_width_px * img_height / img_width)
                    slide.shapes.add_picture(img_path, 0, (slide_height_px - scaled_height) // 2,
                                              width=slide_width_px, height=scaled_height)
                elif fit_mode == "fit_height":
                    scaled_width = int(slide_height_px * img_width / img_height)
                    slide.shapes.add_picture(img_path, (slide_width_px - scaled_width) // 2, 0,
                                              width=scaled_width, height=slide_height_px)
                elif fit_mode == "letterbox":
                    slide.shapes.add_picture(img_path, 0, 0, width=slide_width_px, height=slide_height_px)
                elif fit_mode == "fill":
                    slide.shapes.add_picture(img_path, 0, 0, width=slide_width_px, height=slide_height_px)

    # Save the PowerPoint presentation
    output_pptx_path = os.path.join(temp_dir, "presentation.pptx")
    presentation.save(output_pptx_path)

    # Return the generated file for download
    return send_from_directory(temp_dir, "presentation.pptx", as_attachment=True)

# Main entry point for running the app
if __name__ == "__main__":
    # Run the app on all network interfaces at port 8080
    # Debug mode is disabled for production
    app.run(host="0.0.0.0", port=8080)
