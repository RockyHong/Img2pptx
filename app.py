"""
app.py
A Flask web server that:
  - Accepts a POST request with a .zip file of images
  - Extracts images, converts them to slides in a PPTX (PowerPoint) presentation
  - Returns the PPTX as a file download
  - Also provides a GET request to serve the HTML form (index.html)

HOW IT WORKS:
1. On GET /: We return the index.html file.
2. On POST /: 
   - We read form fields: slide_width_px, slide_height_px, fit_mode
   - We read a .zip of images
   - We unzip images to a temporary folder
   - We create a PPTX, set slide size, add each image to a new slide
   - We optionally add a rectangle shape for debug (you can remove this)
   - We save the PPTX in memory (BytesIO) to avoid file system issues
   - We return that PPTX to the user as a downloadable file
3. We clean up the temporary folder after each operation.

BEST PRACTICES:
- For production, consider setting a maximum file size.
- For large uploads, store files in Google Cloud Storage instead of a local temp.
"""

from flask import Flask, request, send_file, send_from_directory
import os
import zipfile
import re
import shutil
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from io import BytesIO

app = Flask(__name__)

def natural_sort_key(s):
    """
    Sort strings in "natural" order, so "slide2" < "slide10", etc.
    Splits the string by digits vs. non-digits.
    """
    return [int(text) if text.isdigit() else text.lower() 
            for text in re.split(r'(\d+)', s)]

@app.route("/", methods=["GET", "POST"])
def home():
    # 1) If GET: Serve the index.html
    if request.method == "GET":
        # This assumes index.html is in the same directory as app.py
        return send_from_directory(".", "index.html")

    # 2) If POST: Process the form & ZIP
    if request.method == "POST":
        # a) Collect form inputs
        slide_width_px = int(request.form.get("slide_width_px", 1920))
        slide_height_px = int(request.form.get("slide_height_px", 1080))
        fit_mode = request.form.get("fit_mode", "fit_width")

        # b) Get the .zip file
        images_zip = request.files.get("images_zip")
        if not images_zip or not images_zip.filename.endswith(".zip"):
            return "Error: Please upload a valid .zip file containing images.", 400

        # c) Create a temporary folder for extraction
        temp_dir = "temp_upload"
        os.makedirs(temp_dir, exist_ok=True)
        zip_path = os.path.join(temp_dir, images_zip.filename)
        # Save the uploaded ZIP to disk
        images_zip.save(zip_path)

        # d) Extract images from the ZIP
        with zipfile.ZipFile(zip_path, "r") as zip_ref:
            for file in zip_ref.namelist():
                # Extract only image files
                if file.lower().endswith((".png", ".jpg", ".jpeg")):
                    zip_ref.extract(file, temp_dir)

        # e) Create a new PPTX
        presentation = Presentation()
        # Convert px to inches at ~96 DPI
        presentation.slide_width = Inches(slide_width_px / 96.0)
        presentation.slide_height = Inches(slide_height_px / 96.0)

        # Typically layout #6 is a blank slide in the default template
        blank_layout = presentation.slide_layouts[6]

        added_files = 0

        # f) Loop over extracted files in natural order
        for root, _, files in os.walk(temp_dir):
            for file in sorted(files, key=natural_sort_key):
                if file.lower().endswith((".png", ".jpg", ".jpeg")):
                    try:
                        img_path = os.path.join(root, file)

                        # Open the image with PIL to get width/height
                        img = Image.open(img_path)
                        w, h = img.size
                        aspect = w / float(h)

                        # Create a new slide
                        slide = presentation.slides.add_slide(blank_layout)

                        # Slide dimensions (inches)
                        slide_w_in = presentation.slide_width
                        slide_h_in = presentation.slide_height

                        # Fit the image per the chosen fit_mode
                        if fit_mode == "fit_width":
                            final_width = slide_w_in
                            final_height = final_width / aspect
                            # If it doesn't fit in height, shrink
                            if final_height > slide_h_in:
                                final_height = slide_h_in
                                final_width = final_height * aspect
                        elif fit_mode == "fit_height":
                            final_height = slide_h_in
                            final_width = final_height * aspect
                            # If it doesn't fit in width, shrink
                            if final_width > slide_w_in:
                                final_width = slide_w_in
                                final_height = final_width / aspect
                        elif fit_mode == "fill":
                            # Fill entire slide, ignoring aspect
                            final_width = slide_w_in
                            final_height = slide_h_in
                        else:
                            # "letterbox" or default
                            final_width = slide_w_in
                            final_height = final_width / aspect
                            if final_height > slide_h_in:
                                final_height = slide_h_in
                                final_width = final_height * aspect

                        # Center the image on the slide
                        left = (slide_w_in - final_width) / 2
                        top = (slide_h_in - final_height) / 2

                        # Add the picture
                        slide.shapes.add_picture(
                            img_path, left, top, 
                            width=final_width, 
                            height=final_height
                        )
                        added_files += 1

                    except Exception as e:
                        print(f"Error processing {file}: {e}")

        # g) If no images found, return an error
        if added_files == 0:
            shutil.rmtree(temp_dir, ignore_errors=True)
            return "No valid images found.", 400

        # h) Save the PPTX in memory instead of disk
        pptx_buffer = BytesIO()
        presentation.save(pptx_buffer)
        pptx_buffer.seek(0)

        file_size = len(pptx_buffer.getvalue())
        print(f"Slides: {added_files}, PPTX size: {file_size} bytes")

        # i) Clean up the temp folder
        shutil.rmtree(temp_dir, ignore_errors=True)

        # j) Return the PPTX as a downloadable file
        return send_file(
            pptx_buffer,
            as_attachment=True,
            download_name="presentation.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

# If running locally:
if __name__ == "__main__":
    # By default, listens on port 8080
    app.run(host="0.0.0.0", port=8080)
